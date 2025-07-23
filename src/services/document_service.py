"""文档处理服务"""

import logging
import asyncio
import hashlib
from pathlib import Path
from typing import List, Dict, Optional, Any
from datetime import datetime
import aiofiles
import mimetypes

from fastapi import UploadFile
from src.config.settings import get_settings
from src.models.document import Document, DocumentInfo, DocumentStatus
from src.models.text_chunk import TextChunk, ChunkMetadata
from src.services.embedding_service import EmbeddingService
from src.services.vector_service import VectorService
from src.services.bm25_service import BM25Service

logger = logging.getLogger(__name__)
settings = get_settings()


class DocumentService:
    """文档处理服务"""
    
    def __init__(self):
        self.upload_dir = Path(settings.upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        
        # 初始化依赖服务
        self.embedding_service = EmbeddingService()
        self.vector_service = VectorService()
        self.bm25_service = BM25Service()
        
        # 支持的文件类型
        self.supported_types = {
            'application/pdf': '.pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'application/vnd.ms-powerpoint': '.ppt',
            'text/plain': '.txt',
            'text/markdown': '.md'
        }
    
    async def upload_document(self, file: UploadFile, metadata: Optional[Dict] = None) -> DocumentInfo:
        """上传文档"""
        try:
            # 验证文件类型
            if file.content_type not in self.supported_types:
                raise ValueError(f"不支持的文件类型: {file.content_type}")
            
            # 创建文档信息
            doc_info = DocumentInfo(
                filename=file.filename,
                original_filename=file.filename,
                file_size=0,  # 将在保存后更新
                mime_type=file.content_type,
                metadata=metadata or {}
            )
            
            # 生成文件路径
            file_extension = self.supported_types[file.content_type]
            file_path = self.upload_dir / f"{doc_info.id}{file_extension}"
            
            # 保存文件
            content = await file.read()
            async with aiofiles.open(file_path, 'wb') as f:
                await f.write(content)
            
            # 更新文件大小
            doc_info.file_size = len(content)
            doc_info.status = DocumentStatus.PROCESSING
            
            logger.info(f"文档上传成功: {doc_info.filename} ({doc_info.file_size} bytes)")
            
            # 异步处理文档
            asyncio.create_task(self._process_document_async(doc_info, file_path))
            
            return doc_info
            
        except Exception as e:
            logger.error(f"文档上传失败: {str(e)}")
            raise
    
    async def _process_document_async(self, doc_info: DocumentInfo, file_path: Path):
        """异步处理文档"""
        try:
            # 解析文档
            content = await self.parse_document(file_path, doc_info.mime_type)
            
            # 分割文本
            chunks = await self.split_text(content, doc_info)
            
            # 生成嵌入向量
            chunks_with_embeddings = await self.embedding_service.embed_chunks(chunks)
            
            # 存储到向量数据库
            await self.vector_service.store_chunks(chunks_with_embeddings)
            
            # 更新BM25索引
            await self.bm25_service.add_documents(chunks_with_embeddings)
            
            # 更新文档状态
            doc_info.status = DocumentStatus.COMPLETED
            
            logger.info(f"文档处理完成: {doc_info.filename}")
            
        except Exception as e:
            logger.error(f"文档处理失败: {str(e)}")
            doc_info.status = DocumentStatus.FAILED
    
    async def parse_document(self, file_path: Path, mime_type: str) -> str:
        """解析文档内容"""
        try:
            if mime_type == 'text/plain' or mime_type == 'text/markdown':
                return await self._parse_text_file(file_path)
            elif mime_type == 'application/pdf':
                return await self._parse_pdf_file(file_path)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                return await self._parse_word_file(file_path)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                return await self._parse_ppt_file(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {mime_type}")
                
        except Exception as e:
            logger.error(f"文档解析失败: {str(e)}")
            raise
    
    async def _parse_text_file(self, file_path: Path) -> str:
        """解析文本文件"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            return await f.read()
    
    async def _parse_pdf_file(self, file_path: Path) -> str:
        """解析PDF文件"""
        try:
            import PyPDF2
            
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            
            return text.strip()
            
        except ImportError:
            logger.error("PyPDF2 未安装，无法解析PDF文件")
            raise
        except Exception as e:
            logger.error(f"PDF解析失败: {str(e)}")
            raise
    
    async def _parse_word_file(self, file_path: Path) -> str:
        """解析Word文件"""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text.strip()
            
        except ImportError:
            logger.error("python-docx 未安装，无法解析Word文件")
            raise
        except Exception as e:
            logger.error(f"Word文档解析失败: {str(e)}")
            raise
    
    async def _parse_ppt_file(self, file_path: Path) -> str:
        """解析PPT文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
            
        except ImportError:
            logger.error("python-pptx 未安装，无法解析PPT文件")
            raise
        except Exception as e:
            logger.error(f"PPT文档解析失败: {str(e)}")
            raise
    
    async def split_text(self, content: str, doc_info: DocumentInfo) -> List[TextChunk]:
        """分割文本为块"""
        try:
            chunks = []
            chunk_size = settings.chunk_size
            chunk_overlap = settings.chunk_overlap
            
            # 简单的文本分割策略
            words = content.split()
            
            for i in range(0, len(words), chunk_size - chunk_overlap):
                chunk_words = words[i:i + chunk_size]
                chunk_content = " ".join(chunk_words)
                
                if not chunk_content.strip():
                    continue
                
                # 创建文本块
                chunk = TextChunk(
                    document_id=str(doc_info.id),
                    content=chunk_content,
                    chunk_index=len(chunks),
                    start_position=i,
                    end_position=min(i + chunk_size, len(words)),
                    metadata=ChunkMetadata(
                        chunk_id="",  # 将在创建后设置
                        document_title=doc_info.filename,
                        chunk_index=len(chunks),
                        total_chunks=0,  # 将在最后更新
                        page_number=1,
                        section_title="",
                        language="zh",
                        confidence_score=1.0
                    )
                )
                
                # 设置chunk_id
                chunk.metadata.chunk_id = str(chunk.id)
                chunks.append(chunk)
            
            # 更新总块数
            for chunk in chunks:
                chunk.metadata.total_chunks = len(chunks)
            
            logger.info(f"文本分割完成: {len(chunks)} 个块")
            return chunks
            
        except Exception as e:
            logger.error(f"文本分割失败: {str(e)}")
            raise
    
    async def delete_document(self, doc_id: str) -> bool:
        """删除文档"""
        try:
            # 从向量数据库删除
            await self.vector_service.delete_by_document(doc_id)
            
            # 从BM25索引删除
            await self.bm25_service.remove_documents([doc_id])
            
            # 删除文件（如果存在）
            # 这里需要根据实际的文件存储策略来实现
            
            logger.info(f"文档删除成功: {doc_id}")
            return True
            
        except Exception as e:
            logger.error(f"文档删除失败: {str(e)}")
            return False
    
    async def get_document_info(self, doc_id: str) -> Optional[DocumentInfo]:
        """获取文档信息"""
        # 这里需要实现文档信息的持久化存储和检索
        # 暂时返回None，后续需要集成数据库
        return None
    
    async def list_documents(self, limit: int = 50, offset: int = 0) -> List[DocumentInfo]:
        """列出文档"""
        # 这里需要实现文档列表的持久化存储和检索
        # 暂时返回空列表，后续需要集成数据库
        return []
